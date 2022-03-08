# misc.py

from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.shapes.autoshape import Shape
from copy import deepcopy
from math import ceil, floor
from datetime import date

from pathlib import Path
ROOT = Path(__file__).parent.parent.absolute()

from .custom_exceptions import *



# This method reads the arguements passed in to determine how many Rounds and
# 	Tracks there will be.
def determine_parameters(ARGS):

	# default is supposed to be mutually exclusive with rounds and tracks, so if
	#     -d is provided, it cannot be combined with -r or -t.
	# One the other end, -r and -t can be used together or separatly, but cannot
	#     be used with -d.
	if not (ARGS.default or ARGS.rounds or ARGS.tracks):
		raise NoFlagsError
	elif ARGS.default and (ARGS.rounds or ARGS.tracks):
		raise AmbiguousNumRoundsTracksError
	else:

		# If -d was provided, set rounds and tracks to the default values of 5
		#     and 10 respectively.
		if ARGS.default: rounds, tracks = 5, 10

		# Otherwise, check to see if -r and/or -t was/were provided.
		# If -r was provided, set rounds to that number, else set it to 5.
		# If -t was provided, set tracks to that number, else set it to 10.
		else: 
			rounds = ARGS.rounds if ARGS.rounds is not None else 5
			tracks = ARGS.tracks if ARGS.tracks is not None else 10

	challenge_round = True if not ARGS.nochallenge else False

	return rounds, tracks, challenge_round
	
	
# This method is for generating a new slide.
def add_slide(prs, layout, title):
	
	# Create the slide based off of the passed in layout.
	slide = prs.slides.add_slide(layout)
	
	# Title the slide.
	slide.shapes.title.text = title

	return slide


# This method is for adding additional lines of text to a slide
def add_text(slide, text, level=0, bold=False, italic=False, underline=False, append=False):

	# Each slide has some number of items on them called shapes, which
	# 	themselves store some number of placeholders.
	# These placeholders are stored in a dictionary within the slide's shapes
	# 	called placeholders{}.
	# To add text to a slide, we have to access the placeholder stored at key 1.
	# Note that since placeholders is not an array, this is actually accessing 
	# 	KEY 1, NOT INDEX 1.
	# From there, we can access the text_frame, which is what stores the text in
	# 	the shape.
	content_area = slide.shapes.placeholders[1].text_frame

	# text_frames are composed of various Paragraphs.
	# p will be the Paragraph that our text argument will be written into.
	# However, before we can write we have to determine whether the text_frame
	# 	already has content.
	# If it does, we have to call .add_paragraph(), a method belonging to
	# 	text_frame that creates a new Paragraph that can be written to.
	# Otherwise, we should write to the Paragraph already there.
	if not append and content_area.text:
		p = content_area.add_paragraph()
	else:
		p = content_area.paragraphs[len(content_area.paragraphs) - 1]

	
	# Write the text into the space.
	run = p.add_run()
	run.text = text

	# Apply text styling.
	p.level = level
	run.font.bold = bold
	run.font.italic = italic
	if underline: run.font.underline = MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE


# This method is for duplicating a shape and adding it to the slide on which it
# 	appears.
def clone_shape(shape):
	shape_obj = shape.element
	sp_tree = shape_obj.getparent()
	new_sp = deepcopy(shape_obj)
	sp_tree.append(new_sp)
	new_shape = Shape(new_sp, None)
	new_shape.left = shape.left
	new_shape.top = shape.top
	new_shape.width = shape.width
	new_shape.height = shape.height
	return new_shape


# Equation for review_slide() that determines where the Track and Text Box
# 	X-coordinates are.
def determine_x_coord(i, type):

	# The starting value is determined by if the X-coordinate is for the slide's
	# 	Track or Text Box
	start = 6.5 if (type == "track") else 7.5
	
	# The offset is for if the item is supposed to go in the left column or the
	# 	right column.
	offset = 0 if (i % 2 != 0) else 4.5

	return start + offset


# Equation for review_slide() that determines where the Track and Text Box
# 	Y-coordinates are.
def determine_y_coord(n, i):

	# Set how far apart the items will be.
	SPREAD = 1.5

	# 4 will place the item right in the middle of the slide.
	# To shift the item up or down depending on i, the following equation is
	# 	evaluated:
	return 4 - (round(

		# y = m * x + b, where
		# m is -SPREAD
		# x is ceil(i / 2) to account for the left and right columns
		# b is (SPREAD * ceil(n / 2)) to find the base y-intercept
		-SPREAD * ceil(i / 2) + (SPREAD * ceil(n / 2))

		# Next make slight adjustments to the above equation.
		# Subtract the amount required to place the middle items along the
		# 	center of the slide.
		- (SPREAD * floor(n / 4))

		# If ceil(n / 2) % 2 == 0, then there will be an even number of rows, so
		# 	everything has to be futher offset by a small amount.
		- (
			0
			if ceil(n / 2) % 2 != 0 else

			# The offset is either positive or negative depending on if n is a
			# 	multiple of 4.
			((SPREAD / 2) if n % 4 != 0 else -(SPREAD / 2))
		)
	, 3))


# This method is for saving the PowerPoint produced by this script.
def save_presentation(presentation):
	
	# Get the current year and season
	today = date.today()
	year = today.year
	season = "Spring" if today.month < 7 else "Fall"

	# Save the passed in presentation.
	presentation.save(ROOT / f"VGMC {year} {season} Slides.pptx")