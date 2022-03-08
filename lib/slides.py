# slides.py

from .misc import *

from pptx.util import Inches, Pt
from io import FileIO

from pathlib import Path
ROOT = Path(__file__).parent.parent.absolute()

# This method creates the Rules slide.
def rules_slide(prs, rounds, tracks, challenge=False):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[1], "Rules and Scoring")

	# Add this text to the slide if it isn't for the special Challenge Round.
	if not challenge:
		add_text(slide, f"There will be {rounds} Rounds each with {tracks} Tracks.")
		add_text(slide, "You will get roughly 30 – 45 seconds of music to guess from.")
		add_text(slide, "Scoring is as follows:")
		add_text(slide, "1 point for Game Franchise", level=1)
		add_text(slide, "1 point for Specific Game", level=1)
		add_text(slide, "1 point for Track Name/Place", level=1)
		add_text(slide, "Some songs don’t have official releases, or play in multiple games/areas. Those songs will have a star listed on the answer key, so if you put something close to the listing you’ll still receive the point.")

	# Else, change the title and add the special Challenge Round rules.
	else:
		slide.shapes.title.text = "Special Rules"
		add_text(slide, "This will be the hardest round.")
		add_text(slide, "This time you only need the game’s name, with a correct guess earning ")
		add_text(slide, "5 points", bold=True, underline=True, append=True)
		add_text(slide, ".", append=True)
		add_text(slide, "Because of this, we are going to be much stricter about getting the name right.")
		add_text(slide, "All songs chosen here are ones on my personal playlists, but might not be the most popular from its associated game.")
		add_text(slide, "Good luck!")

# This method creates the slides that show the round number before each round.
def round_slide(prs, round_num, challenge=False):

	# Create and title the slide.
	# Because slide_layouts[2] doesn't have a title placeholder, the slide
	# 	creation and titling has to be explicitly done rather than with the
	# 	usual add_slide().
	slide = prs.slides.add_slide(prs.slide_layouts[2])
	tf = slide.shapes.placeholders[10].text_frame
	tf.paragraphs[0].text = f"Round {round_num}"
	
	# If the round is supposed to be a challenge round, this extra bit of
	# 	formatted text also needs be added to the slide's "title" field.
	if challenge:
		p = tf.add_paragraph()
		p.text = "(Challenge Round)"
		p.font.size = Pt(66)


# This method creates the individual Track slides for each round.
def track_slide(prs, round_num, track_num):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[3], f"Track {track_num}")

	# Add the Audio object.
	slide.shapes.add_movie(
		FileIO(ROOT / f"tracks/{round_num}-{track_num}.mp3", "rb"),
		left=Inches(8), top=Inches(1.5), width=Inches(6), height=Inches(6),
		poster_frame_image=FileIO(ROOT / "templates" / "speaker.png", "rb"),
		mime_type='audio/mp3',
	)


# This method creates the Review slide for relistening to every track in the
# 	round.
def review_slide(prs, round_num, num_tracks):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[4], f"Round {round_num} Review")

	# Create a reference to the base text placeholder that will be copied.
	template_text_box = slide.shapes.placeholders[10]


	# Generate the song-text pairs.
	for index in range(1, num_tracks + 1):

		# Add the track.
		track = slide.shapes.add_movie(

			# Location of the track.
			FileIO(ROOT / "tracks" / f"{round_num}-{index}.mp3", "rb"),

			# X-coordinate of the track
			left=Inches(determine_x_coord(index, "track")),
			
			# Y-coordinate of the track
			top=Inches(determine_y_coord(num_tracks, index)),
			
			# Dimensions of the track
			width=Inches(1), height=Inches(1),

			# Extra details
			poster_frame_image=FileIO(ROOT / "templates" / "speaker.png", "rb"),
			mime_type='audio/mp3'
		)


		# Add the text box.
		# Start by cloning the template text placeholder.
		cloned_tb = clone_shape(template_text_box)

		# Set the X-coordinate.
		cloned_tb.left = Inches(determine_x_coord(index, "text"))

		# Set the Y-coordinate.
		cloned_tb.top = Inches(determine_y_coord(num_tracks, index))

		# Add the text.
		cloned_tb.text = f" Track {index}"

		# If num_tracks is odd, the final Track and Text Box are sitting alone
		# 	in the left column with no corresponding item in the right column,
		# 	so horizontally center it.
		if (index == num_tracks) and (num_tracks % 2 != 0):
			track.left = Inches(8.75)
			cloned_tb.left = Inches(9.75)
		

	# Finally delete the template text box.
	sp = template_text_box._sp
	sp.getparent().remove(sp)


# This method creates the Answer slide for showing the round answers.
def answer_slide(prs, answers_ws, round_num, num_tracks, challenge=False):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[5], f"Round {round_num} Answers")

	# Because the answers are stored sequentially, and have a variable number of
	# 	tracks, an offset needs to be generated.
	# An additional +1 is also added to account for the first row being headers
	# 	in the Excel Worksheet.
	offset = ((round_num - 1) * num_tracks) + 1

	# Add the answers to the slide.
	for i in range(1, num_tracks + 1):

		# If it's not a challenge round, add the Game Franchise.
		if not challenge:
			add_text(slide, f"{answers_ws[f'A{offset + i}'].value}: ")

		# Next add the Game.
		# The append status is dependent on if challenge is still False.
		# If so, it should just append to the Franchise paragraph.
		# If not, there is no Franchise paragraph so it should NOT append.
		add_text(slide, f"{answers_ws[f'B{offset + i}'].value}", italic=True, append=(not challenge))
		
		# Finally add the Song.
		add_text(slide, f" – {answers_ws[f'C{offset + i}'].value}", append=True)
