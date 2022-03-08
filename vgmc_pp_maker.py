# vgmc_pp_maker.py
from argparse import ArgumentParser
from copy import deepcopy
from datetime import date
from io import FileIO
from math import ceil, floor
from pathlib import Path
from sys import exit
ROOT = Path(__file__).parent.absolute()

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from openpyxl import load_workbook


# Parse the arguments provided to argparse.
PARSER = ArgumentParser(description="Use this file to create a PowerPoint for the Florida Southern College Computer Science Club's Video Game Music Competiton (FSCCSCVGMC).")
PARSER.add_argument('-d', '--default', action="store_true", help='creates a default PowerPoint with 5 Rounds each with 10 Tracks.')
PARSER.add_argument('-r', '--rounds', type=int, nargs='?', help='provide a number of rounds')
PARSER.add_argument('-t', '--tracks', type=int, nargs='?', help='provide a number of tracks for each round')
PARSER.add_argument('-nc', '--nochallenge', action='store_true', help="use this flag if you don't want to have the final round be a challege round")
ARGS = PARSER.parse_args()

def main():
	
	# Determine how many Rounds and Tracks per Round there will be.
	num_rounds, num_tracks, challenge = determine_parameters()
	print(f"\033[96mThere will be \033[93m{num_rounds}\033[96m Round{'s Each' if num_rounds > 1 else ''} with \033[93m{num_tracks}\033[96m Tracks.\033[0m")
	print(f"\033[96mRound {num_rounds} will {'also' if challenge else 'not'} be a Challenge Round.\033[0m")

	# Create the PowerPoint using the slide template.
	prs = Presentation(ROOT / "templates" / "master_copy.pptx")

	# Read in the answer key for later
	wb = load_workbook(ROOT / 'tracks' / 'song_info.xlsx')
	answer_sheet = wb.active

	# Create the Title and Rules Slides
	rules_slide(prs, num_rounds, num_tracks)

	# Create the slides for each round.
	for round_num in range(1, num_rounds + 1):

		challenge_round = (challenge and round_num == num_rounds)

		round_slide(prs, round_num, challenge_round)

		if challenge_round:
			rules_slide(prs, num_rounds, num_tracks, True)

		for track_num in range(1, num_tracks + 1):
			track_slide(prs, round_num, track_num)

		review_slide(prs, round_num, num_tracks)

		answer_slide(prs, answer_sheet, round_num, num_tracks, challenge_round)


	# Save the Presentation
	save_presentation(prs)


# This method reads the arguements passed in to determine how many Rounds and
# 	Tracks there will be.
def determine_parameters():

	# default is supposed to be mutually exclusive with rounds and tracks, so if
	#     -d is provided, it cannot be combined with -r or -t.
	# One the other end, -r and -t can be used together or separatly, but cannot
	#     be used with -d.
	if not (ARGS.default):
		print("You have to use flags to tell the program how many Rounds and Tracks you want.")
		print("Try using the -h flag to learn more.")
		exit(2)	
	elif ARGS.default and (ARGS.rounds or ARGS.tracks):
		print("--default and --rounds|--tracks are mutually exclusive and cannot be used together.")
		exit(2)
	else:

		# If -d was provided, set rounds and tracks to the default values of 5
		#     and 10 respectively.
		if ARGS.default: rounds, tracks = 5, 10

		# Otherwise, check to see if -r and/or -t was/were provided.
		# If -r was provided, set rounds to that number, else set it to 5.
		# If -t was provided, set tracks to that number, else set it to 10.
		else: 
			rounds = ARGS.rounds if ARGS.rounds is not None else 5
			tracks = ARGS.tracks if ARGS.tracks is not None else 10

	challenge_round = True if not ARGS.nochallenge else False

	return rounds, tracks, challenge_round


# This method is for generating a new slide.
def add_slide(prs, layout, title):
	
	# Create the slide based off of the passed in layout.
	slide = prs.slides.add_slide(layout)
	
	# Title the slide.
	slide.shapes.title.text = title

	return slide


# This method is for adding additional lines of text to a slide
def add_text(slide, text, level=0, bold=False, italic=False, underline=False, append=False):

	# Each slide has some number of items on them called shapes, which
	# 	themselves store some number of placeholders.
	# These placeholders are stored in a dictionary within the slide's shapes
	# 	called placeholders{}.
	# To add text to a slide, we have to access the placeholder stored at key 1.
	# Note that since placeholders is not an array, this is actually accessing 
	# 	KEY 1, NOT INDEX 1.
	# From there, we can access the text_frame, which is what stores the text in
	# 	the shape.
	content_area = slide.shapes.placeholders[1].text_frame

	# text_frames are composed of various Paragraphs.
	# p will be the Paragraph that our text argument will be written into.
	# However, before we can write we have to determine whether the text_frame
	# 	already has content.
	# If it does, we have to call .add_paragraph(), a method belonging to
	# 	text_frame that creates a new Paragraph that can be written to.
	# Otherwise, we should write to the Paragraph already there.
	if not append and content_area.text:
		p = content_area.add_paragraph()
	else:
		p = content_area.paragraphs[len(content_area.paragraphs) - 1]

	
	# Write the text into the space.
	run = p.add_run()
	run.text = text

	# Apply text styling.
	p.level = level
	run.font.bold = bold
	run.font.italic = italic
	if underline: run.font.underline = MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE


# This method is for duplicating a shape and adding it to the slide on which it
# 	appears.
def clone_shape(shape):
	shape_obj = shape.element
	sp_tree = shape_obj.getparent()
	new_sp = deepcopy(shape_obj)
	sp_tree.append(new_sp)
	new_shape = Shape(new_sp, None)
	new_shape.left = shape.left
	new_shape.top = shape.top
	new_shape.width = shape.width
	new_shape.height = shape.height
	return new_shape


# This method creates the Rules slide.
def rules_slide(prs, rounds, tracks, challenge=False):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[1], "Rules and Scoring")

	# Add this text to the slide if it isn't for the special Challenge Round.
	if not challenge:
		add_text(slide, f"There will be {rounds} Rounds each with {tracks} Tracks.")
		add_text(slide, "You will get roughly 30 – 45 seconds of music to guess from.")
		add_text(slide, "Scoring is as follows:")
		add_text(slide, "1 point for Game Franchise", level=1)
		add_text(slide, "1 point for Specific Game", level=1)
		add_text(slide, "1 point for Track Name/Place", level=1)
		add_text(slide, "Some songs don’t have official releases, or play in multiple games/areas. Those songs will have a star listed on the answer key, so if you put something close to the listing you’ll still receive the point.")

	# Else, change the title and add the special Challenge Round rules.
	else:
		slide.shapes.title.text = "Special Rules"
		add_text(slide, "This will be the hardest round.")
		add_text(slide, "This time you only need the game’s name, with a correct guess earning ")
		add_text(slide, "5 points", bold=True, underline=True, append=True)
		add_text(slide, ".", append=True)
		add_text(slide, "Because of this, we are going to be much stricter about getting the name right.")
		add_text(slide, "All songs chosen here are ones on my personal playlists, but might not be the most popular from its associated game.")
		add_text(slide, "Good luck!")

# This method creates the slides that show the round number before each round.
def round_slide(prs, round_num, challenge=False):

	# Create and title the slide.
	# Because slide_layouts[2] doesn't have a title placeholder, the slide
	# 	creation and titling has to be explicitly done rather than with the
	# 	usual add_slide().
	slide = prs.slides.add_slide(prs.slide_layouts[2])
	tf = slide.shapes.placeholders[10].text_frame
	tf.paragraphs[0].text = f"Round {round_num}"
	
	# If the round is supposed to be a challenge round, this extra bit of
	# 	formatted text also needs be added to the slide's "title" field.
	if challenge:
		p = tf.add_paragraph()
		p.text = "(Challenge Round)"
		p.font.size = Pt(66)


# This method creates the individual Track slides for each round.
def track_slide(prs, round_num, track_num):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[3], f"Track {track_num}")

	# Add the Audio object.
	slide.shapes.add_movie(
		FileIO(ROOT / f"tracks/{round_num}-{track_num}.mp3", "rb"),
		left=Inches(8), top=Inches(1.5), width=Inches(6), height=Inches(6),
		poster_frame_image=FileIO(ROOT / "templates" / "speaker.png", "rb"),
		mime_type='audio/mp3',
	)


# This method creates the Review slide for relistening to every track in the
# 	round.
def review_slide(prs, round_num, num_tracks):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[4], f"Round {round_num} Review")

	# Create a reference to the base text placeholder that will be copied.
	template_text_box = slide.shapes.placeholders[10]


	# Generate the song-text pairs.
	for index in range(1, num_tracks + 1):

		# Add the track.
		track = slide.shapes.add_movie(

			# Location of the track.
			FileIO(ROOT / "tracks" / f"{round_num}-{index}.mp3", "rb"),

			# X-coordinate of the track
			left=Inches(determine_x_coord(index, "track")),
			
			# Y-coordinate of the track
			top=Inches(determine_y_coord(num_tracks, index)),
			
			# Dimensions of the track
			width=Inches(1), height=Inches(1),

			# Extra details
			poster_frame_image=FileIO(ROOT / "templates" / "speaker.png", "rb"),
			mime_type='audio/mp3'
		)


		# Add the text box.
		# Start by cloning the template text placeholder.
		cloned_tb = clone_shape(template_text_box)

		# Set the X-coordinate.
		cloned_tb.left = Inches(determine_x_coord(index, "text"))

		# Set the Y-coordinate.
		cloned_tb.top = Inches(determine_y_coord(num_tracks, index))

		# Add the text.
		cloned_tb.text = f" Track {index}"

		# If num_tracks is odd, the final Track and Text Box are sitting alone
		# 	in the left column with no corresponding item in the right column,
		# 	so horizontally center it.
		if (index == num_tracks) and (num_tracks % 2 != 0):
			track.left = Inches(8.75)
			cloned_tb.left = Inches(9.75)
		

	# Finally delete the template text box.
	sp = template_text_box._sp
	sp.getparent().remove(sp)


# Equation for review_slide() that determines where the Track and Text Box
# 	X-coordinates are.
def determine_x_coord(i, type):

	# The starting value is determined by if the X-coordinate is for the slide's
	# 	Track or Text Box
	start = 6.5 if (type == "track") else 7.5
	
	# The offset is for if the item is supposed to go in the left column or the
	# 	right column.
	offset = 0 if (i % 2 != 0) else 4.5

	return start + offset


# Equation for review_slide() that determines where the Track and Text Box
# 	Y-coordinates are.
def determine_y_coord(n, i):

	# Set how far apart the items will be.
	SPREAD = 1.5

	# 4 will place the item right in the middle of the slide.
	# To shift the item up or down depending on i, the following equation is
	# 	evaluated:
	return 4 - (round(

		# y = m * x + b, where
		# m is -SPREAD
		# x is ceil(i / 2) to account for the left and right columns
		# b is (SPREAD * ceil(n / 2)) to find the base y-intercept
		-SPREAD * ceil(i / 2) + (SPREAD * ceil(n / 2))

		# Next make slight adjustments to the above equation.
		# Subtract the amount required to place the middle items along the
		# 	center of the slide.
		- (SPREAD * floor(n / 4))

		# If ceil(n / 2) % 2 == 0, then there will be an even number of rows, so
		# 	everything has to be futher offset by a small amount.
		- (
			0
			if ceil(n / 2) % 2 != 0 else

			# The offset is either positive or negative depending on if n is a
			# 	multiple of 4.
			((SPREAD / 2) if n % 4 != 0 else -(SPREAD / 2))
		)
	, 3))


# This method creates the Answer slide for showing the round answers.
def answer_slide(prs, answers_ws, round_num, num_tracks, challenge=False):

	# Create and title the slide.
	slide = add_slide(prs, prs.slide_layouts[5], f"Round {round_num} Answers")

	# Because the answers are stored sequentially, and have a variable number of
	# 	tracks, an offset needs to be generated.
	# An additional +1 is also added to account for the first row being headers
	# 	in the Excel Worksheet.
	offset = ((round_num - 1) * num_tracks) + 1

	# Add the answers to the slide.
	for i in range(1, num_tracks + 1):

		# If it's not a challenge round, add the Game Franchise.
		if not challenge:
			add_text(slide, f"{answers_ws[f'A{offset + i}'].value}: ")

		# Next add the Game.
		# The append status is dependent on if challenge is still False.
		# If so, it should just append to the Franchise paragraph.
		# If not, there is no Franchise paragraph so it should NOT append.
		add_text(slide, f"{answers_ws[f'B{offset + i}'].value}", italic=True, append=(not challenge))
		
		# Finally add the Song.
		add_text(slide, f" – {answers_ws[f'C{offset + i}'].value}", append=True)


# This method is for saving the PowerPoint produced by this script.
def save_presentation(presentation):
	
	# Get the current year and season
	today = date.today()
	year = today.year
	season = "Spring" if today.month < 7 else "Fall"

	# Save the passed in presentation.
	presentation.save(ROOT / f"VGMC {year} {season} Slides.pptx")


if __name__ == '__main__':
	print('\n\033[92mRunning vgmc_pp_maker.py\033[0m')

	try:
		main()

	except Exception as e:
		print(f'\033[91m{e}\033[0m')